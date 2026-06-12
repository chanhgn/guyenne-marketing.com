#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère un deck PowerPoint 16:9 éditable pour la présentation KOBO
(Portrait dirigeante — Sandrine Cruchon), aux couleurs de la marque.
Reprend la direction artistique de la vidéo Remotion : anthracite + cuivre,
logo, cartes, photos.  Sortie : out/KOBO-Presentation.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PUB = os.path.join(ROOT, "public")
OUT = os.path.join(ROOT, "out", "KOBO-Presentation.pptx")

# ---- Palette de marque ----
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_CARD = RGBColor(0x18, 0x24, 0x3E)
INK_CARD2 = RGBColor(0x14, 0x1F, 0x36)
COPPER = RGBColor(0xCA, 0x8A, 0x04)
COPPER_GLOW = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
SLATE_FAINT = RGBColor(0x64, 0x74, 0x8B)
HAIRLINE = RGBColor(0x2A, 0x3A, 0x5C)
FONT = "Inter"

EMU_IN = 914400
SW, SH = 13.333, 7.5  # pouces


def img(name):
    return os.path.join(PUB, name)


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, w=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def no_shadow(shape):
    # python-pptx ajoute parfois une ombre par défaut sur les formes auto
    sp = shape._element.spPr
    existing = sp.find(qn("a:effectLst"))
    if existing is None:
        sp.append(sp.makeelement(qn("a:effectLst"), {}))


def bg(slide, color=INK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, x, y, w, h, color, rounded=False, ln=None, lnw=0.75):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    solid(shp, color)
    if ln is None:
        no_line(shp)
    else:
        line(shp, ln, lnw)
    no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def corner_bracket(slide, size=1.0, x=None, y=0.45, thick=0.028):
    """Équerre cuivre en haut à droite."""
    if x is None:
        x = SW - 0.45 - size
    h = rect(slide, x, y, size, thick, COPPER_GLOW)
    v = rect(slide, x + size - thick, y, thick, size, COPPER_GLOW)
    return h, v


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines = list of dicts: {text, size, color, bold, italic, spacing(after pt),
    caps(letter-tracking via spc), line_spacing}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        if "space_before" in ln:
            p.space_before = Pt(ln["space_before"])
        if "line_spacing" in ln:
            p.line_spacing = ln["line_spacing"]
        # plusieurs runs possible (pour mêler couleurs)
        runs = ln["runs"] if "runs" in ln else [ln]
        for r in runs:
            run = p.add_run()
            run.text = r["text"]
            run.font.name = FONT
            run.font.size = Pt(r.get("size", 16))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            run.font.color.rgb = r.get("color", WHITE)
            if r.get("tracking"):
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(r["tracking"] * 100)))
    return tb


def kicker(slide, text, x, y):
    rect(slide, x, y + 0.085, 0.32, 0.03, COPPER_GLOW)
    w = min(9.0, SW - (x + 0.42) - 0.35)
    textbox(slide, x + 0.42, y - 0.06, w, 0.4, [
        {"text": text.upper(), "size": 12.5, "bold": True, "color": COPPER_GLOW, "tracking": 3},
    ])


def title(slide, text, x, y, w=11.0, size=34):
    textbox(slide, x, y, w, 1.4, [
        {"text": text, "size": size, "bold": True, "color": WHITE, "line_spacing": 1.02},
    ])


def add_cover(slide, path, x, y, w, h, radius=False):
    iw, ih = Image.open(path).size
    img_ar = iw / ih
    box_ar = w / h
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if img_ar > box_ar:
        crop = (1 - box_ar / img_ar) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ar / box_ar) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.line.color.rgb = HAIRLINE
    pic.line.width = Pt(1)
    return pic


def add_logo(slide, x, y, w, full=True):
    name = "kobo_logo_light_full.png" if full else "kobo_logo_light.png"
    p = img(name)
    iw, ih = Image.open(p).size
    h = w * ih / iw
    slide.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    return h


def chip(slide, path, x, y, w, h):
    """Logo partenaire sur pastille blanche."""
    rect(slide, x, y, w, h, WHITE, rounded=True, ln=HAIRLINE, lnw=0.75)
    pad = 0.16
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = w - 2 * pad, h - 2 * pad
    if ar > bw / bh:
        pw = bw
        ph = bw / ar
    else:
        ph = bh
        pw = bh * ar
    slide.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                             Inches(pw), Inches(ph))


def card(slide, x, y, w, h, num, head, body, org=None):
    c = rect(slide, x, y, w, h, INK_CARD, rounded=True, ln=HAIRLINE, lnw=0.75)
    rect(slide, x, y + 0.12, 0.06, h - 0.24, COPPER)  # accent gauche
    tb = slide.shapes.add_textbox(Inches(x + 0.42), Inches(y + 0.34), Inches(w - 0.8), Inches(h - 0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    # numéro
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = num
    r0.font.name = FONT; r0.font.size = Pt(26); r0.font.bold = True
    r0.font.color.rgb = COPPER
    p0.space_after = Pt(6)
    # titre
    p1 = tf.add_paragraph()
    r1 = p1.add_run(); r1.text = head
    r1.font.name = FONT; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = WHITE
    p1.space_after = Pt(4)
    if org:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = org.upper()
        r2.font.name = FONT; r2.font.size = Pt(11); r2.font.bold = True
        r2.font.color.rgb = COPPER_GLOW
        rPr = r2._r.get_or_add_rPr(); rPr.set("spc", "150")
        p2.space_after = Pt(8)
    p3 = tf.add_paragraph()
    r3 = p3.add_run(); r3.text = body
    r3.font.name = FONT; r3.font.size = Pt(13); r3.font.color.rgb = SLATE
    p3.line_spacing = 1.2


def pillar(slide, x, y, w, h, num, head, body, head_size=15, body_size=12):
    rect(slide, x, y, w, h, INK_CARD, rounded=True, ln=HAIRLINE, lnw=0.75)
    # badge numéro
    b = rect(slide, x + 0.3, y + 0.32, 0.55, 0.55, INK_CARD2, rounded=True, ln=COPPER, lnw=1)
    textbox(slide, x + 0.3, y + 0.42, 0.55, 0.4,
            [{"text": num, "size": 15, "bold": True, "color": COPPER_GLOW, "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(x + 0.32), Inches(y + 1.05), Inches(w - 0.62), Inches(h - 1.25))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = head
    r1.font.name = FONT; r1.font.size = Pt(head_size); r1.font.bold = True; r1.font.color.rgb = WHITE
    p1.space_after = Pt(6); p1.line_spacing = 1.05
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = body
    r2.font.name = FONT; r2.font.size = Pt(body_size); r2.font.color.rgb = SLATE
    p2.line_spacing = 1.18


def band(slide, x, y, w, label, text, h=0.8):
    rect(slide, x, y, w, h, INK_CARD2, rounded=True, ln=COPPER, lnw=1)
    rect(slide, x + 0.28, y + h / 2 - 0.06, 0.12, 0.12, COPPER_GLOW, rounded=True)
    tb = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y), Inches(w - 0.9), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if label:
        rl = p.add_run(); rl.text = label.upper() + "   "
        rl.font.name = FONT; rl.font.size = Pt(12); rl.font.bold = True; rl.font.color.rgb = COPPER_GLOW
        rPr = rl._r.get_or_add_rPr(); rPr.set("spc", "150")
    rt = p.add_run(); rt.text = text
    rt.font.name = FONT; rt.font.size = Pt(13.5); rt.font.color.rgb = WHITE
    p.line_spacing = 1.12


TOTAL_SLIDES = 13


def page_no(slide, n):
    textbox(slide, SW - 1.6, SH - 0.55, 1.2, 0.3,
            [{"text": f"{n:02d} / {TOTAL_SLIDES:02d}", "size": 11, "bold": True,
              "color": SLATE_FAINT, "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)


def base(prs, n, with_bracket=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, INK)
    if with_bracket:
        corner_bracket(s)
    if n:
        page_no(s, n)
    return s


# =========================================================================
def build(out_path=OUT):
    global TOTAL_SLIDES
    TOTAL_SLIDES = 13
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    MX = 0.85  # marge gauche

    # ---- 1. Couverture ----
    s = base(prs, None, with_bracket=False)
    corner_bracket(s, size=1.1, x=SW - 0.55 - 1.1, y=0.5)
    # équerre bas-gauche (orientation correcte : ouvre vers le haut-droite)
    _bx, _by, _bs, _th = 0.55, SH - 0.5, 1.1, 0.028
    rect(s, _bx, _by - _th, _bs, _th, COPPER_GLOW)       # barre horizontale (bas)
    rect(s, _bx, _by - _bs, _th, _bs, COPPER_GLOW)       # barre verticale (montante)
    lh = add_logo(s, SW / 2 - 2.6, 2.1, 5.2, full=True)
    rect(s, SW / 2 - 1.6, 2.1 + lh + 0.35, 3.2, 0.03, COPPER_GLOW)
    textbox(s, 1, 2.1 + lh + 0.6, SW - 2, 1.0, [
        {"text": "Sandrine Cruchon", "size": 26, "bold": True, "color": WHITE,
         "align": PP_ALIGN.CENTER, "space_after": 6},
        {"text": "GÉRANTE ASSOCIÉE — KOBO MENUISERIE", "size": 13, "bold": True,
         "color": SLATE, "align": PP_ALIGN.CENTER, "tracking": 2.5},
    ], align=PP_ALIGN.CENTER)

    # ---- 2. Portrait ----
    s = base(prs, 2)
    add_cover(s, img("kobo_portrait_sandrine.png"), MX, 1.0, 4.6, 5.5, radius=True)
    tx = MX + 5.1
    kicker(s, "Portrait dirigeante", tx, 1.05)
    title(s, "Une gérante passionnée", tx, 1.5, 7.0, 34)
    textbox(s, tx, 2.85, 6.9, 2.6, [
        {"text": "Née à Bordeaux en juin 1978, j’ai construit ma trajectoire autour de deux "
                 "piliers : une rigueur technique absolue, issue du terrain, et un profond "
                 "dévouement familial.", "size": 14.5, "color": SLATE, "line_spacing": 1.3,
         "space_after": 12},
        {"text": "Mère de deux enfants — un aîné de 19 ans en études supérieures, un cadet de "
                 "15 ans au lycée.", "size": 14.5, "color": SLATE, "line_spacing": 1.3},
    ])
    rect(s, tx, 5.35, 0.05, 1.0, COPPER_GLOW)
    textbox(s, tx + 0.25, 5.4, 6.6, 1.0, [
        {"text": "Une vie dédiée à l’art des façades et de l’enveloppe du bâtiment.",
         "size": 15, "italic": True, "color": WHITE, "line_spacing": 1.25},
    ])

    # ---- 3. Fondations académiques ----
    s = base(prs, 3)
    kicker(s, "Fondations académiques", MX, 0.95)
    title(s, "Une expertise née de l’ingénierie", MX, 1.4, 11, 32)
    cy, ch2, cw = 2.85, 3.6, 5.55
    card(s, MX, cy, cw, ch2, "01", "BTS Enveloppe du Bâtiment",
         "Formation hautement technique : physique du bâtiment, structures porteuses, "
         "thermique des menuiseries, étude globale de la peau des édifices.",
         org="Études supérieures à Anglet")
    card(s, MX + cw + 0.45, cy, cw, ch2, "02", "DNTS Spécialiste des Façades",
         "Spécialisation de pointe : menuiserie aluminium, façades-rideaux complexes et "
         "vitrages structurels — théorie d’ingénierie et pratique terrain.",
         org="Formation à Évry, en alternance")

    # ---- 4. Premiers grands chantiers ----
    s = base(prs, 4)
    kicker(s, "Premiers grands chantiers", MX, 0.85)
    title(s, "L’exigence du terrain", MX, 1.3, 11, 32)
    cw4 = 4.35
    card(s, MX, 2.55, cw4, 2.95, "01", "Apprentissage chez REALCO",
         "Un premier poste fondateur pour s’imprégner de l’exigence opérationnelle et de la "
         "réalité des chantiers de menuiserie.")
    card(s, MX + cw4 + 0.4, 2.55, cw4, 2.95, "02", "Conduite de travaux chez CANCE",
         "Pilotage financier et technique de chantiers tertiaires d’envergure, de 25 000 € à "
         "plus de 3 000 000 €.")
    px = MX + 2 * cw4 + 0.8
    add_cover(s, img("kobo_icade.png"), px, 2.55, 2.95, 1.4)
    add_cover(s, img("kobo_immeuble_toulouse.png"), px, 4.1, 2.95, 1.4)
    band(s, MX, 5.75, 11.6, "Référence",
         "Supervision de l’enveloppe architecturale du complexe tertiaire ICADE à Toulouse.")

    # ---- 5. L'essor d'une entrepreneuse ----
    s = base(prs, 5)
    kicker(s, "L’essor d’une entrepreneuse", MX, 0.95)
    title(s, "De la conduite de travaux à la reprise", MX, 1.4, 11, 30)
    cw5 = 4.35
    card(s, MX, 2.85, cw5, 3.3, "01", "Alkar Aluminium",
         "Conduite de chantiers complexes à forte valeur technique. Chantiers majeurs : "
         "l’Hôpital de Mauléon et la façade iconique du Mama Shelter à Toulouse.")
    card(s, MX + cw5 + 0.4, 2.85, cw5, 3.3, "02", "Reprise de MDG Fermetures",
         "Une opportunité d’entreprendre née de l’amour du métier : reprise d’une PME locale "
         "au départ en retraite du dirigeant historique. MDG devient KOBO.")
    add_cover(s, img("kobo_mauleon.png"), MX + 2 * cw5 + 0.8, 2.85, 2.95, 3.3)

    # ---- 6. Une PME proche de vous ----
    s = base(prs, 6)
    kicker(s, "Une PME proche de vous", MX, 0.85)
    title(s, "Gages de confiance", MX, 1.3, 11, 32)
    pw, ph, gap = 3.7, 2.75, 0.3
    items6 = [
        ("01", "Proximité & écoute", "PME à taille humaine basée à Tournefeuille, en synergie "
         "directe avec les propriétaires et acteurs tertiaires de l’agglomération toulousaine."),
        ("02", "Spécialisation technique", "Focus absolu sur la menuiserie haut de gamme, "
         "l’aluminium technique sur-mesure et l’optimisation thermique des bâtiments."),
        ("03", "Certifiée RGE & Qualibat", "La certitude d’un travail d’expert, ouvrant droit "
         "aux aides de l’État pour l’efficacité énergétique."),
    ]
    for i, (n, hd, bd) in enumerate(items6):
        pillar(s, MX + i * (pw + gap), 2.55, pw, ph, n, hd, bd)
    chip(s, img("kobo_qualibat.png"), MX, 5.55, 1.5, 0.95)
    chip(s, img("kobo_synerciel.png"), MX + 1.7, 5.55, 1.9, 0.95)

    # ---- 7. Rigueur & précision ----
    s = base(prs, 7)
    kicker(s, "Rigueur & précision KOBO", MX, 0.95)
    title(s, "Un process sans surprise", MX, 1.4, 11, 32)
    cw7 = 5.55
    card(s, MX, 2.85, cw7, 3.4, "1", "Conseil & étude clairs",
         "Chaque projet démarre par un diagnostic sur place. Des outils d’ingénierie de pointe "
         "(Pro Devis) pour des chiffrages rigoureux, transparents, détaillés poste par poste.")
    card(s, MX + cw7 + 0.45, 2.85, cw7, 3.4, "2", "Relevé de côtes sur site",
         "Une commande d’usine sans erreur exige une précision infaillible : le relevé de côtes "
         "final pour la fabrication est systématiquement réalisé en personne.")

    # ---- 8. La qualité de pose directe ----
    s = base(prs, 8)
    kicker(s, "La qualité de pose directe", MX, 0.75)
    title(s, "Des partenaires d’élite, un suivi dirigeant", MX, 1.2, 11.5, 30)
    pw8 = 3.7
    items8 = [
        ("01", "Partenaires d'élite", "Fabrication confiée à des fournisseurs français et "
         "européens : Fenêtréa, Méo, Ates, Bel’m et la Maison Cadiou pour les portails."),
        ("02", "4 équipes de spécialistes", "Installateurs qualifiés, segmentés par corps de "
         "métier : rénovation, grands projets tertiaires, clôtures & portails, portes de garage."),
        ("03", "Suivi 100 % dirigeant", "Aucun intermédiaire technique. Conformité esthétique et "
         "étanchéité pilotées en direct par Sandrine ou Jean-François."),
    ]
    for i, (n, hd, bd) in enumerate(items8):
        pillar(s, MX + i * (pw8 + 0.3), 2.4, pw8, 2.75, n, hd, bd)
    partners = ["kobo_fenetrea.png", "kobo_schuco.png", "kobo_meo.png", "kobo_cadiou.png", "kobo_ates.png"]
    cwp = 2.05
    for i, p in enumerate(partners):
        chip(s, img(p), MX + i * (cwp + 0.2), 5.5, cwp, 1.0)

    # ---- 9. Étude de cas : Villa de Balma ----
    s = base(prs, 9)
    add_cover(s, img("kobo_villa_balma.jpg"), MX, 1.0, 5.0, 5.2, radius=True)
    tx = MX + 5.5
    kicker(s, "Étude de cas — Villa de Balma", tx, 1.0)
    title(s, "Rénovation complète", tx, 1.45, 6.6, 30)
    bullets9 = [
        ("Le déclencheur", "Projet confié à la suite de la recommandation enthousiaste d’un client conquis."),
        ("Le défi", "Mixer le PVC plaxé à l’étage et l’aluminium Fenêtréa haute isolation au RDC pour ajuster le budget."),
        ("L'astuce KOBO", "Préservation éco-responsable des persiennes d’origine — sablage lourd et thermolaquage RAL 7016 Anthracite."),
    ]
    by = 2.7
    for lead, body in bullets9:
        rect(s, tx, by + 0.06, 0.14, 0.14, COPPER_GLOW, rounded=True)
        textbox(s, tx + 0.32, by - 0.04, 6.3, 1.0, [
            {"runs": [
                {"text": lead + " — ", "size": 13.5, "bold": True, "color": COPPER_GLOW},
                {"text": body, "size": 13.5, "color": SLATE},
            ], "line_spacing": 1.2},
        ])
        by += 1.0
    band(s, tx, 5.75, 6.6, "Bénéfice",
         "Harmonie visuelle, confort thermique idéal et budget optimisé.")

    # ---- 10. Étude de cas : Mairie de Rabastens ----
    s = base(prs, 10)
    add_cover(s, img("kobo_mairie_rabastens.jpg"), SW - MX - 5.0, 1.0, 5.0, 5.2, radius=True)
    tx = MX
    kicker(s, "Étude de cas tertiaire — Mairie de Rabastens", tx, 1.0)
    title(s, "Rénovation complète", tx, 1.45, 6.6, 30)
    bullets10 = [
        ("Le déclencheur", "Projet remporté à la suite d’un appel d’offres public."),
        ("Le défi", "Secteur ABF (Architecte des Bâtiments de France) et aluminium Schüco haute isolation, cintré."),
    ]
    by = 2.8
    for lead, body in bullets10:
        rect(s, tx, by + 0.06, 0.14, 0.14, COPPER_GLOW, rounded=True)
        textbox(s, tx + 0.32, by - 0.04, 6.2, 1.0, [
            {"runs": [
                {"text": lead + " — ", "size": 13.5, "bold": True, "color": COPPER_GLOW},
                {"text": body, "size": 13.5, "color": SLATE},
            ], "line_spacing": 1.2},
        ])
        by += 1.1
    band(s, tx, 5.75, 6.5, "Bénéfice",
         "Harmonie visuelle, confort thermique idéal et budget optimisé.")

    # ---- 11. Pourquoi choisir KOBO ----
    s = base(prs, 11)
    kicker(s, "Pourquoi choisir KOBO ?", MX, 0.75)
    title(s, "Quatre raisons décisives", MX, 1.2, 11, 32)
    pw11 = 2.75
    items11 = [
        ("01", "Technicité", "Fabrication sur-mesure en interne : souplesse, maîtrise des délais et qualité constante."),
        ("02", "Showroom de proximité", "Un espace à Tournefeuille pour apprécier en réel les finitions et la robustesse."),
        ("03", "Circuit court", "Aucun commercial tiers : dialogue direct avec les deux gérants-associés."),
        ("04", "Hybride B2C & B2B", "Méthodologies adaptées aux villas des particuliers comme aux structures tertiaires."),
    ]
    for i, (n, hd, bd) in enumerate(items11):
        pillar(s, MX + i * (pw11 + 0.27), 2.45, pw11, 2.9, n, hd, bd, head_size=14, body_size=11.5)
    band(s, MX, 5.7, 11.6, "Financement",
         "Solutions de paiement échelonné via notre partenariat avec la Société Générale (France Finances).")

    # ---- 12. La confiance par les faits ----
    s = base(prs, 12)
    kicker(s, "La confiance par les faits", MX, 0.85)
    title(s, "Notre gage de confiance", MX, 1.3, 11, 32)
    stats = [("4.9", "/5", "Satisfaction client"),
             ("50", "ans", "d'expérience cumulée"),
             ("100", "+/an", "chantiers uniques")]
    sw_, sh_, sgap = 3.7, 2.2, 0.3
    for i, (val, unit, lab) in enumerate(stats):
        x = MX + i * (sw_ + sgap)
        rect(s, x, 2.6, sw_, sh_, INK_CARD, rounded=True, ln=HAIRLINE, lnw=0.75)
        tb = slide_stat = s.shapes.add_textbox(Inches(x), Inches(2.85), Inches(sw_), Inches(1.2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rv = p.add_run(); rv.text = val
        rv.font.name = FONT; rv.font.size = Pt(58); rv.font.bold = True; rv.font.color.rgb = WHITE
        ru = p.add_run(); ru.text = " " + unit
        ru.font.name = FONT; ru.font.size = Pt(24); ru.font.bold = True; ru.font.color.rgb = COPPER_GLOW
        textbox(s, x, 4.15, sw_, 0.5, [
            {"text": lab, "size": 14, "color": SLATE, "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)
    textbox(s, MX, 5.25, 11.6, 1.6, [
        {"text": "En alliant les parcours de Sandrine Cruchon et Jean-François Mores, KOBO réunit "
                 "plus de 50 ans d’expérience dans la fermeture et les façades. Chaque année, plus "
                 "de 100 chantiers uniques prennent vie sous notre supervision exclusive à travers "
                 "l’Occitanie, portés par des certifications d’État reconnues (Qualibat, RGE).",
         "size": 13.5, "color": SLATE, "line_spacing": 1.3}])

    # ---- 13. Merci / contact ----
    s = base(prs, None, with_bracket=False)
    corner_bracket(s, size=1.1, x=SW - 0.55 - 1.1, y=0.5)
    lh = add_logo(s, SW / 2 - 1.85, 1.5, 3.7, full=True)
    textbox(s, 1, 1.5 + lh + 0.5, SW - 2, 1.0, [
        {"text": "Merci pour votre attention", "size": 30, "bold": True, "color": WHITE,
         "align": PP_ALIGN.CENTER, "space_after": 6},
        {"text": "Des questions ?", "size": 18, "bold": True, "color": COPPER_GLOW,
         "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.6, 4.7, 3.2, 0.03, COPPER_GLOW)
    textbox(s, 1, 4.95, SW - 2, 1.2, [
        {"text": "contact@kobo-alu.fr     |     05 61 85 51 40     |     www.kobo-alu.fr",
         "size": 15, "color": SLATE, "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "164 Chemin de Larramet, 31170 Tournefeuille", "size": 12.5,
         "color": SLATE_FAINT, "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print("Saved:", out_path, "—", len(prs.slides._sldIdLst), "slides")


# =========================================================================
# Helpers spécifiques au deck BNI (orienté recommandation)
# =========================================================================
def tag_card(slide, x, y, w, h, text):
    rect(slide, x, y, w, h, INK_CARD, rounded=True, ln=HAIRLINE, lnw=0.75)
    rect(slide, x + w / 2 - 0.18, y + 0.24, 0.36, 0.045, COPPER_GLOW)
    textbox(slide, x + 0.1, y + 0.1, w - 0.2, h - 0.2,
            [{"text": text, "size": 12.5, "bold": True, "color": WHITE,
              "align": PP_ALIGN.CENTER, "line_spacing": 1.05}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def quote_card(slide, x, y, w, h, text):
    rect(slide, x, y, w, h, INK_CARD, rounded=True, ln=HAIRLINE, lnw=0.75)
    rect(slide, x, y + 0.12, 0.06, h - 0.24, COPPER)
    textbox(slide, x + 0.32, y, w - 0.55, h,
            [{"text": "« " + text + " »", "size": 13.5, "italic": True, "color": WHITE,
              "line_spacing": 1.2}], anchor=MSO_ANCHOR.MIDDLE)


def check_row(slide, x, y, w, lead, rest):
    textbox(slide, x, y - 0.02, 0.4, 0.4,
            [{"text": "✓", "size": 17, "bold": True, "color": COPPER_GLOW}])
    textbox(slide, x + 0.42, y, w - 0.42, 0.9,
            [{"runs": [
                {"text": lead + " ", "size": 14.5, "bold": True, "color": WHITE},
                {"text": rest, "size": 14.5, "color": SLATE},
            ], "line_spacing": 1.2}])


def build_bni(out_path):
    """Deck BNI : pensé pour le PUBLIC (comment recommander KOBO)."""
    global TOTAL_SLIDES
    TOTAL_SLIDES = 10
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    MX = 0.85

    # ---- 1. Couverture ----
    s = base(prs, None, with_bracket=False)
    corner_bracket(s, size=1.1, x=SW - 0.55 - 1.1, y=0.5)
    _bx, _by, _bs, _th = 0.55, SH - 0.5, 1.1, 0.028
    rect(s, _bx, _by - _th, _bs, _th, COPPER_GLOW)
    rect(s, _bx, _by - _bs, _th, _bs, COPPER_GLOW)
    lh = add_logo(s, SW / 2 - 2.6, 1.95, 5.2, full=True)
    rect(s, SW / 2 - 1.6, 1.95 + lh + 0.32, 3.2, 0.03, COPPER_GLOW)
    textbox(s, 1, 1.95 + lh + 0.55, SW - 2, 1.4, [
        {"text": "Sandrine Cruchon", "size": 25, "bold": True, "color": WHITE,
         "align": PP_ALIGN.CENTER, "space_after": 5},
        {"text": "GÉRANTE ASSOCIÉE — KOBO MENUISERIE", "size": 12.5, "bold": True,
         "color": SLATE, "align": PP_ALIGN.CENTER, "tracking": 2.5, "space_after": 14},
        {"text": "Menuiseries & aménagements extérieurs  ·  Tournefeuille – Toulouse",
         "size": 14, "color": COPPER_GLOW, "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER)

    # ---- 2. Ce que je fais, simplement ----
    s = base(prs, 2)
    kicker(s, "KOBO en 10 secondes", MX, 0.85)
    title(s, "Ce que je fais, simplement", MX, 1.3, 11, 32)
    textbox(s, MX, 2.7, 11.6, 0.9, [
        {"text": "Je remplace et je pose tout ce qui ouvre et ferme votre maison.",
         "size": 21, "bold": True, "color": WHITE, "line_spacing": 1.1}])
    cats = ["Fenêtres\n& baies vitrées", "Portes\n& portails", "Volets\n& stores",
            "Pergolas\n& carports", "Garde-corps\n& clôtures"]
    cw = (11.633 - 4 * 0.25) / 5
    for i, c in enumerate(cats):
        tag_card(s, MX + i * (cw + 0.25), 3.95, cw, 1.25, c)
    band(s, MX, 5.7, 11.633, "Du conseil à la pose",
         "Aluminium · PVC · Bois — sur-mesure, neuf & rénovation, sur Tournefeuille, "
         "Toulouse & environs.")

    # ---- 3. En qui vous avez confiance (express) ----
    s = base(prs, 3)
    add_cover(s, img("kobo_portrait_sandrine.png"), MX, 1.0, 4.4, 5.5, radius=True)
    tx = MX + 4.9
    kicker(s, "À qui vous avez affaire", tx, 1.05)
    title(s, "Une experte de l’enveloppe du bâtiment", tx, 1.5, 7.2, 27)
    trust = [
        ("Ingénieure des façades —", "BTS Enveloppe du Bâtiment + DNTS spécialiste façades & aluminium."),
        ("+ de 50 ans d’expérience cumulée —", "avec mon associé Jean-François Mores."),
        ("Certifiée RGE & Qualibat —", "vos clients ont droit aux aides de l’État pour l’énergie."),
        ("4.9/5 de satisfaction —", "plus de 100 chantiers uniques par an en Occitanie."),
        ("Fabrication interne & suivi dirigeant —", "pas de commercial : vous parlez au patron."),
    ]
    yy = 2.95
    for lead, rest in trust:
        check_row(s, tx, yy, 7.2, lead, rest)
        yy += 0.78

    # ---- 4. Qui je cherche à rencontrer (client idéal) ----
    s = base(prs, 4)
    kicker(s, "Mon client idéal", MX, 0.85)
    title(s, "À qui penser pour me recommander", MX, 1.3, 11.5, 31)
    pw = 2.75
    targets = [
        ("01", "Particuliers", "Propriétaires de maison qui rénovent ou font construire (Toulouse Ouest et agglomération)."),
        ("02", "Prescripteurs", "Architectes, maîtres d’œuvre, constructeurs de maisons individuelles."),
        ("03", "Immobilier", "Agences immobilières, syndics de copropriété, gestionnaires de biens."),
        ("04", "Tertiaire & public", "Mairies, collectivités, entreprises — façades, secteur ABF, appels d’offres."),
    ]
    for i, (n, hd, bd) in enumerate(targets):
        pillar(s, MX + i * (pw + 0.27), 2.5, pw, 3.4, n, hd, bd, head_size=16, body_size=12)
    band(s, MX, 6.05, 11.633, "Astuce",
         "Pas besoin que la personne sache ce qu’elle veut : si elle a une maison et un projet, pensez à moi.")

    # ---- 5. Les phrases déclic ----
    s = base(prs, 5)
    kicker(s, "Le réflexe recommandation", MX, 0.85)
    title(s, "Si vous entendez ça, pensez à moi", MX, 1.3, 11.5, 31)
    phrases = [
        "Mes fenêtres sont vieilles, j’ai froid et de la condensation.",
        "J’aimerais une pergola ou un carport pour cet été.",
        "Mon portail ou mon volet roulant est en panne.",
        "Je rénove ma maison / je fais construire.",
        "Je veux profiter des aides RGE pour mieux isoler.",
        "Notre copropriété ou notre mairie a un projet de façade.",
    ]
    qw, qh = 5.6, 1.18
    for i, ph in enumerate(phrases):
        cx = MX + (i % 2) * (qw + 0.4)
        cy = 2.45 + (i // 2) * (qh + 0.22)
        quote_card(s, cx, cy, qw, qh, ph)

    # ---- 6. Recommandez-moi en confiance ----
    s = base(prs, 6)
    kicker(s, "Zéro risque pour vous", MX, 0.85)
    title(s, "Recommandez-moi en confiance", MX, 1.3, 11.5, 31)
    pw6 = 3.7
    derisk = [
        ("01", "Fabrication interne", "Sur-mesure fabriqué en interne : délais maîtrisés et qualité constante."),
        ("02", "Suivi 100 % dirigeant", "Sandrine ou Jean-François contrôlent chaque chantier en direct."),
        ("03", "Partenaires d’élite", "Fenêtréa, Schüco, Méo, Maison Cadiou, Ates — et un showroom à Tournefeuille."),
    ]
    for i, (n, hd, bd) in enumerate(derisk):
        pillar(s, MX + i * (pw6 + 0.3), 2.45, pw6, 2.6, n, hd, bd)
    chip(s, img("kobo_qualibat.png"), MX, 5.35, 1.45, 0.95)
    chip(s, img("kobo_fenetrea.png"), MX + 1.65, 5.35, 2.0, 0.95)
    chip(s, img("kobo_schuco.png"), MX + 3.85, 5.35, 1.6, 0.95)
    chip(s, img("kobo_cadiou.png"), MX + 5.65, 5.35, 2.0, 0.95)
    textbox(s, MX + 7.95, 5.35, 3.6, 0.95, [
        {"text": "4.9/5", "size": 26, "bold": True, "color": WHITE, "space_after": 0},
        {"text": "100+ chantiers / an · RGE", "size": 12, "color": SLATE}],
        anchor=MSO_ANCHOR.MIDDLE)

    # ---- 7. Ils nous ont fait confiance (preuves) ----
    s = base(prs, 7)
    kicker(s, "La preuve par les réalisations", MX, 0.85)
    title(s, "Ils nous ont fait confiance", MX, 1.3, 11.5, 31)
    cwp = 5.6
    add_cover(s, img("kobo_villa_balma.jpg"), MX, 2.5, cwp, 2.7)
    add_cover(s, img("kobo_mairie_rabastens.jpg"), MX + cwp + 0.4, 2.5, cwp, 2.7)
    textbox(s, MX, 5.35, cwp, 1.4, [
        {"text": "Villa de Balma — particulier", "size": 14, "bold": True, "color": COPPER_GLOW,
         "space_after": 3},
        {"text": "Rénovation complète alu + PVC, persiennes d’origine restaurées. "
                 "Projet né d’une recommandation client.", "size": 12.5, "color": SLATE,
         "line_spacing": 1.2}])
    textbox(s, MX + cwp + 0.4, 5.35, cwp, 1.4, [
        {"text": "Mairie de Rabastens — tertiaire / ABF", "size": 14, "bold": True,
         "color": COPPER_GLOW, "space_after": 3},
        {"text": "Façade aluminium Schüco cintrée, secteur classé. Remportée sur appel d’offres.",
         "size": 12.5, "color": SLATE, "line_spacing": 1.2}])

    # ---- 8. Ma demande cette semaine (LE slide clé) ----
    s = base(prs, 8)
    kicker(s, "Le plus important", MX, 0.9)
    title(s, "Ma demande cette semaine", MX, 1.4, 11, 34)
    rect(s, MX, 2.95, 11.633, 1.7, INK_CARD, rounded=True, ln=COPPER, lnw=1.25)
    rect(s, MX, 3.1, 0.08, 1.4, COPPER)
    textbox(s, MX + 0.5, 2.95, 10.8, 1.7, [
        {"runs": [
            {"text": "Présentez-moi à ", "size": 22, "color": WHITE},
            {"text": "une seule personne", "size": 22, "bold": True, "color": COPPER_GLOW},
            {"text": " qui change ses fenêtres, rénove sa maison ou veut une pergola.",
             "size": 22, "color": WHITE},
        ], "line_spacing": 1.18}], anchor=MSO_ANCHOR.MIDDLE)
    band(s, MX, 5.05, 11.633, "Cette semaine en particulier",
         "Je cherche à rencontrer un architecte ou un constructeur de maisons individuelles "
         "de l’agglomération toulousaine.", h=1.0)
    textbox(s, MX, 6.25, 11.633, 0.5, [
        {"text": "→ une recommandation précise vaut dix contacts vagues.",
         "size": 13, "italic": True, "color": SLATE_FAINT}])

    # ---- 9. Comment me recommander ----
    s = base(prs, 9)
    kicker(s, "Passez-moi le relais", MX, 0.85)
    title(s, "Comment me recommander", MX, 1.3, 11, 32)
    pw9 = 3.7
    steps = [
        ("1", "Donnez mon nom", "Transmettez mon nom et mon numéro à la personne concernée."),
        ("2", "Ou ma carte / le showroom", "Remettez-lui ma carte ou invitez-la à passer au showroom de Tournefeuille."),
        ("3", "Le mieux : en direct", "Mettez-nous directement en relation — c’est le plus efficace pour conclure."),
    ]
    for i, (n, hd, bd) in enumerate(steps):
        pillar(s, MX + i * (pw9 + 0.3), 2.55, pw9, 2.7, n, hd, bd)
    band(s, MX, 5.7, 11.633, "Donnant-donnant",
         "Une bonne recommandation = un projet bien fait et un client qui VOUS remerciera.")

    # ---- 10. Merci / contact ----
    s = base(prs, None, with_bracket=False)
    corner_bracket(s, size=1.1, x=SW - 0.55 - 1.1, y=0.5)
    _bx, _by, _bs, _th = 0.55, SH - 0.5, 1.1, 0.028
    rect(s, _bx, _by - _th, _bs, _th, COPPER_GLOW)
    rect(s, _bx, _by - _bs, _th, _bs, COPPER_GLOW)
    lh = add_logo(s, SW / 2 - 1.85, 1.45, 3.7, full=True)
    textbox(s, 1, 1.45 + lh + 0.45, SW - 2, 1.2, [
        {"text": "Merci — et pensez à KOBO !", "size": 30, "bold": True, "color": WHITE,
         "align": PP_ALIGN.CENTER, "space_after": 6},
        {"text": "Créateur d’ouvertures", "size": 18, "italic": True, "color": COPPER_GLOW,
         "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.6, 4.75, 3.2, 0.03, COPPER_GLOW)
    textbox(s, 1, 5.0, SW - 2, 1.2, [
        {"text": "contact@kobo-alu.fr     |     05 61 85 51 40     |     www.kobo-alu.fr",
         "size": 15, "color": SLATE, "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "164 Chemin de Larramet, 31170 Tournefeuille", "size": 12.5,
         "color": SLATE_FAINT, "align": PP_ALIGN.CENTER},
    ], align=PP_ALIGN.CENTER)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print("Saved:", out_path, "—", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    build_bni(os.path.join(ROOT, "out", "KOBO-Presentation-BNI.pptx"))
    build(os.path.join(ROOT, "out", "KOBO-Portrait-Dirigeante.pptx"))
