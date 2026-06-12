#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Assemble les slides "deck éditorial" KOBO (PNG 1920x1080 rendus via Remotion,
composition KoboDeck) en un PowerPoint 16:9 plein cadre.
Entrée : out/deck_png/slide-01.png .. slide-10.png
Sortie : out/KOBO-BNI-Sexy.pptx
"""
import os
import glob
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(ROOT, "out", "deck_png")
OUT = os.path.join(ROOT, "out", "KOBO-BNI-Sexy.pptx")
EMU_IN = 914400
SW, SH = 13.333, 7.5

P_NS = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'


def add_fade(slide, spd="med"):
    """Ajoute une transition en fondu (convertie en animation par Google Slides)."""
    sld = slide._element
    xml = f'<p:transition {P_NS} spd="{spd}"><p:fade/></p:transition>'
    trans = parse_xml(xml)
    clr = sld.find(qn("p:clrMapOvr"))
    if clr is not None:
        clr.addnext(trans)
    else:
        sld.find(qn("p:cSld")).addnext(trans)


def build():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    blank = prs.slide_layouts[6]
    pngs = sorted(glob.glob(os.path.join(SRC, "slide-*.png")))
    if not pngs:
        raise SystemExit("Aucune image dans " + SRC)
    for p in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
        add_fade(s)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("Saved:", OUT, "—", len(pngs), "slides (transitions fondu)")


if __name__ == "__main__":
    build()
